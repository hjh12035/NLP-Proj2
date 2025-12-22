import os
import shutil
from typing import List, Dict, Optional
import pytesseract
from PIL import Image
import fitz
import docx2txt
from PyPDF2 import PdfReader
from pptx import Presentation

from config import DATA_DIR


class DocumentLoader:
    MIN_IMAGE_SIDE = 512  # 最小图片边长，单位像素

    @staticmethod
    def _is_small_bitmap(width: Optional[int], height: Optional[int]) -> bool:
        if not width or not height:
            return False
        return (
            width < DocumentLoader.MIN_IMAGE_SIDE
            or height < DocumentLoader.MIN_IMAGE_SIDE
        )
    
    def __init__(
        self,
        data_dir: str = DATA_DIR,
    ):
        self.data_dir = data_dir
        self.supported_formats = [".pdf", ".pptx", ".docx", ".txt"]
        self.image_formats = ["png", "jpg", "jpeg", "bmp"]

    def load_pdf(self, file_path: str) -> List[Dict]:
        """加载PDF文件，按页返回内容

        TODO: 实现PDF文件加载
        要求：
        1. 使用PdfReader读取PDF文件
        2. 遍历每一页，提取文本内容
        3. 格式化为"--- 第 X 页 ---\n文本内容\n"
        4. 返回pdf内容列表，每个元素包含 {"text": "..."}
        """
        try:
            reader = PdfReader(file_path)
            pages = []
            for page_num, page in enumerate(reader.pages):
                text = page.extract_text() or ""
                pages.append({"text": f"--- 第 {page_num + 1} 页 ---\n{text}\n"})
            return pages
        except Exception as e:
            print(f"加载PDF文件失败: {file_path}, 错误: {e}")
            return []

    def load_pptx(self, file_path: str) -> List[Dict]:
        """加载PPT文件，按幻灯片返回内容

        TODO: 实现PPT文件加载
        要求：
        1. 使用Presentation读取PPT文件
        2. 遍历每一页，提取文本内容
        3. 格式化为"--- 幻灯片 X ---\n文本内容\n"
        4. 返回幻灯片内容列表，每个元素包含 {"text": "..."}
        """
        try:
            presentation = Presentation(file_path)
            slides = []
            for slide_num, slide in enumerate(presentation.slides):
                slide_text = []
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        slide_text.append(shape.text)
                text = "\n".join(slide_text)
                slides.append({"text": f"--- 幻灯片 {slide_num + 1} ---\n{text}\n"})
            return slides
        except Exception as e:
            print(f"加载PPT文件失败: {file_path}, 错误: {e}")
            return []

    def load_docx(self, file_path: str) -> str:
        """加载DOCX文件
        TODO: 实现DOCX文件加载
        要求：
        1. 使用docx2txt读取DOCX文件
        2. 返回文本内容
        """
        try:
            text = docx2txt.process(file_path)
            return text
        except Exception as e:
            print(f"加载DOCX文件失败: {file_path}, 错误: {e}")
            return ""

    def load_txt(self, file_path: str) -> str:
        """加载TXT文件
        TODO: 实现TXT文件加载
        要求：
        1. 使用open读取TXT文件（注意使用encoding="utf-8"）
        2. 返回文本内容
        """
        try:
            with open(file_path, "r", encoding="utf-8") as f:
                text = f.read()
            return text
        except Exception as e:
            print(f"加载TXT文件失败: {file_path}, 错误: {e}")
            return ""

    def ocr_image(self, image_path: str) -> str:
        """对图片进行OCR识别，返回文本"""
        try:
            text = pytesseract.image_to_string(Image.open(image_path), lang="chi_sim+eng")
            return text.strip()
        except Exception as e:
            print(f"OCR识别失败: {image_path}, 错误: {e}")
            return ""
        
    def extract_images_from_pdf(self, file_path: str, output_dir: str) -> List[Dict]:
        """提取PDF中的图片并进行OCR识别
        
        1. 使用fitz库打开PDF文件
        2. 遍历每一页，提取图片
        3. 保存图片到output_dir
        4. 对每张图片进行OCR识别，获取文本
        5. 格式化为"--- 第 X 页 ---\n图片ocr内容\n"
        6. 返回图片内容列表，每个元素包含 {"filepath": "...", "text": "..."}
        """
        if not os.path.exists(output_dir):
            os.makedirs(output_dir, exist_ok=True)
        filename = os.path.basename(file_path).split(".pdf")[0]
        doc = fitz.open(file_path)
        images = []

        for page_idx in range(len(doc)):
            page = doc[page_idx]
            image_list = page.get_images(full=True)

            for img_idx, img in enumerate(image_list):
                xref = img[0]
                base_image = doc.extract_image(xref)

                width = base_image.get("width")
                height = base_image.get("height")
                if self._is_small_bitmap(width, height) or base_image["ext"].lower() not in self.image_formats:
                    continue

                image_bytes = base_image["image"]
                image_ext = base_image["ext"]
                image_filename = f"pdf{filename}_p{page_idx}_img{img_idx}.{image_ext}"
                image_path = os.path.join(output_dir, image_filename)
                with open(image_path, "wb") as f:
                    f.write(image_bytes)
                ocr_text = self.ocr_image(image_path)
                if not ocr_text:
                    continue
                text = f"--- 第 {page_idx + 1} 页 ---\n[图片OCR内容]\n{ocr_text}\n"
                images.append(
                    {
                        "filepath": image_path,
                        "text": text,
                        "page_number": page_idx + 1,
                        "image_id": img_idx,
                    }
                )
        return images
    
    def extract_images_from_pptx(self, file_path: str, output_dir: str) -> List[Dict]:
        """提取PPT中的图片并进行OCR识别
        1. 使用Presentation读取PPT文件
        2. 遍历每一页，提取图片
        3. 保存图片到output_dir
        4. 对每张图片进行OCR识别，获取文本
        5. 格式化为"--- 幻灯片 X ---\n图片ocr内容\n
        6. 返回图片内容列表，每个元素包含 {"filepath": "...", "text": "..."}
        """
        if not os.path.exists(output_dir):
            os.makedirs(output_dir, exist_ok=True)
        filename = os.path.basename(file_path).split(".ppt")[0]
        prs = Presentation(file_path)
        images = []

        for slide_idx, slide in enumerate(prs.slides):
            for shape_idx, shape in enumerate(slide.shapes):
                if shape.shape_type == 13:  # PICTURE
                    width_px = int(shape.width / 9525)
                    height_px = int(shape.height / 9525)
                    if self._is_small_bitmap(width_px, height_px) or shape.image.ext.lower() not in self.image_formats:
                        continue
                    
                    image = shape.image
                    image_ext = image.ext
                    image_bytes = image.blob

                    image_filename = f"ppt{filename}_s{slide_idx}_img{shape_idx}.{image_ext}"
                    image_path = os.path.join(output_dir, image_filename)

                    with open(image_path, "wb") as f:
                        f.write(image_bytes)

                    ocr_text = self.ocr_image(image_path)
                    if not ocr_text:
                        continue
                    text = f"--- 幻灯片 {slide_idx + 1} ---\n[图片OCR内容]\n{ocr_text}\n"

                    images.append(
                        {
                            "filepath": image_path,
                            "text": text,
                            "page_number": slide_idx + 1,
                            "image_id": shape_idx,
                        }
                    )
        return images


    def load_document(self, file_path: str, image_output_dir: Optional[str] = None) -> List[Dict[str, str]]:
        """加载单个文档，PDF和PPT按页/幻灯片分割，返回文档块列表"""
        ext = os.path.splitext(file_path)[1].lower()
        filename = os.path.basename(file_path)
        documents = []

        if ext == ".pdf":
            pages = self.load_pdf(file_path)
            for page_idx, page_data in enumerate(pages, 1):
                documents.append(
                    {
                        "content": page_data["text"],
                        "filename": filename,
                        "filepath": file_path,
                        "filetype": ext,
                        "page_number": page_idx,
                        "chunk_type": "text",
                    }
                )
            if image_output_dir:
                images = self.extract_images_from_pdf(file_path, image_output_dir)
                if images:
                    for img in images:
                        documents.append(
                            {
                                "content": img.get("text", ""),
                                "filename": filename,
                                "filepath": img["filepath"],
                                "filetype": ext,
                                "page_number": img["page_number"],
                                "image_id": img.get("image_id", 0),
                                "chunk_type": "image",
                            }
                        )
        elif ext == ".pptx":
            slides = self.load_pptx(file_path)
            for slide_idx, slide_data in enumerate(slides, 1):
                documents.append(
                    {
                        "content": slide_data["text"],
                        "filename": filename,
                        "filepath": file_path,
                        "filetype": ext,
                        "page_number": slide_idx,
                        "chunk_type": "text",
                    }
                )
            if image_output_dir:
                images = self.extract_images_from_pptx(file_path, image_output_dir)
                if images:
                    for img in images:
                        documents.append(
                            {
                                "content": img.get("text", ""),
                                "filename": filename,
                                "filepath": img["filepath"],
                                "filetype": ext,
                                "page_number": img["page_number"],
                                "image_id": img.get("image_id", 0),
                                "chunk_type": "image",
                            }
                        )
        elif ext == ".docx":
            content = self.load_docx(file_path)
            if content:
                documents.append(
                    {
                        "content": content,
                        "filename": filename,
                        "filepath": file_path,
                        "filetype": ext,
                        "page_number": 0,
                        "chunk_type": "text",
                    }
                )
        elif ext == ".txt":
            content = self.load_txt(file_path)
            if content:
                documents.append(
                    {
                        "content": content,
                        "filename": filename,
                        "filepath": file_path,
                        "filetype": ext,
                        "page_number": 0,
                        "chunk_type": "text",
                    }
                )
        else:
            print(f"不支持的文件格式: {ext}")

        return documents

    def load_all_documents(self) -> List[Dict[str, str]]:
        """加载数据目录下的所有文档"""
        image_output_dir = os.path.join(self.data_dir, "images")
        if not os.path.exists(self.data_dir):
            print(f"数据目录不存在: {self.data_dir}")
            return None
        if os.path.exists(image_output_dir):
            shutil.rmtree(image_output_dir)
            print("图片目录已清空")
        os.makedirs(image_output_dir, exist_ok=True)

        documents = []

        for root, dirs, files in os.walk(self.data_dir):
            for file in files:
                ext = os.path.splitext(file)[1].lower()
                if ext in self.supported_formats:
                    file_path = os.path.join(root, file)
                    print(f"正在加载: {file_path}")
                    doc_chunks = self.load_document(file_path, image_output_dir=image_output_dir)
                    if doc_chunks:
                        documents.extend(doc_chunks)

        return documents
